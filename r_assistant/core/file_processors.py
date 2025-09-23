"""
文件处理器 - 用于处理各种格式的上传文件
"""
import os
import csv
import json
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Dict, Any, List
import logging

try:
    import pandas as pd
    from docx import Document
    from pptx import Presentation
    from openpyxl import load_workbook
    OFFICE_SUPPORT = True
except ImportError:
    OFFICE_SUPPORT = False

logger = logging.getLogger(__name__)


class FileProcessor:
    """文件处理器基类"""
    
    @staticmethod
    def get_file_info(file) -> Dict[str, Any]:
        """获取文件基本信息"""
        file_path = Path(file.name) if hasattr(file, 'name') else Path(str(file))
        
        return {
            'filename': file_path.name,
            'extension': file_path.suffix.lower(),
            'size': getattr(file, 'size', 0),
            'content': '',
            'preview': '',
            'file_type': FileProcessor.get_file_type(file_path.suffix.lower())
        }
    
    @staticmethod
    def get_file_type(extension: str) -> str:
        """根据文件扩展名确定文件类型"""
        type_mapping = {
            '.txt': 'text',
            '.csv': 'csv',
            '.json': 'json',
            '.xml': 'xml',
            '.xlsx': 'excel',
            '.xls': 'excel',
            '.docx': 'word',
            '.doc': 'word',
            '.pptx': 'powerpoint',
            '.ppt': 'powerpoint',
            '.r': 'r_script',
            '.rmd': 'r_markdown',
            '.py': 'python',
            '.js': 'javascript',
            '.html': 'html',
            '.css': 'css'
        }
        return type_mapping.get(extension, 'unknown')
    
    @staticmethod
    def process_file(file) -> Dict[str, Any]:
        """处理文件并提取内容"""
        file_info = FileProcessor.get_file_info(file)
        extension = file_info['extension']
        
        try:
            if extension == '.txt':
                file_info.update(FileProcessor.process_text_file(file))
            elif extension == '.csv':
                file_info.update(FileProcessor.process_csv_file(file))
            elif extension == '.json':
                file_info.update(FileProcessor.process_json_file(file))
            elif extension == '.xml':
                file_info.update(FileProcessor.process_xml_file(file))
            elif extension in ['.xlsx', '.xls'] and OFFICE_SUPPORT:
                file_info.update(FileProcessor.process_excel_file(file))
            elif extension in ['.docx', '.doc'] and OFFICE_SUPPORT:
                file_info.update(FileProcessor.process_word_file(file))
            elif extension in ['.pptx', '.ppt'] and OFFICE_SUPPORT:
                file_info.update(FileProcessor.process_powerpoint_file(file))
            elif extension in ['.r', '.rmd', '.py', '.js', '.html', '.css']:
                file_info.update(FileProcessor.process_code_file(file))
            else:
                file_info['content'] = f"不支持的文件格式: {extension}"
                file_info['preview'] = "无法预览此文件格式"
                
        except Exception as e:
            logger.error(f"处理文件 {file_info['filename']} 时出错: {str(e)}")
            file_info['content'] = f"文件处理出错: {str(e)}"
            file_info['preview'] = "文件处理失败"
        
        return file_info
    
    @staticmethod
    def process_text_file(file) -> Dict[str, str]:
        """处理文本文件"""
        try:
            content = file.read().decode('utf-8')
            preview = content[:500] + "..." if len(content) > 500 else content
            return {
                'content': content,
                'preview': preview
            }
        except UnicodeDecodeError:
            # 尝试其他编码
            file.seek(0)
            try:
                content = file.read().decode('gbk')
                preview = content[:500] + "..." if len(content) > 500 else content
                return {
                    'content': content,
                    'preview': preview
                }
            except:
                return {
                    'content': "无法解码文件内容",
                    'preview': "编码错误"
                }
    
    @staticmethod
    def process_csv_file(file) -> Dict[str, str]:
        """处理CSV文件"""
        try:
            content = file.read().decode('utf-8')
            file.seek(0)
            
            # 分析CSV结构
            dialect = csv.Sniffer().sniff(content[:1024])
            reader = csv.reader(content.splitlines(), dialect)
            rows = list(reader)
            
            if len(rows) > 0:
                headers = rows[0]
                sample_rows = rows[1:6]  # 前5行数据作为预览
                
                preview = f"CSV文件 ({len(rows)-1} 行数据)\n"
                preview += f"列名: {', '.join(headers)}\n\n"
                preview += "数据预览:\n"
                for i, row in enumerate(sample_rows):
                    if i < 3:  # 只显示前3行
                        preview += f"第{i+1}行: {', '.join(row[:5])}\n"  # 只显示前5列
                
                if len(rows) > 6:
                    preview += "..."
                    
                return {
                    'content': content,
                    'preview': preview,
                    'rows': len(rows) - 1,
                    'columns': len(headers) if headers else 0
                }
            else:
                return {
                    'content': content,
                    'preview': "空的CSV文件"
                }
                
        except Exception as e:
            return {
                'content': f"CSV处理错误: {str(e)}",
                'preview': "CSV解析失败"
            }
    
    @staticmethod
    def process_json_file(file) -> Dict[str, str]:
        """处理JSON文件"""
        try:
            content = file.read().decode('utf-8')
            json_data = json.loads(content)
            
            # 创建结构化预览
            preview = f"JSON文件结构:\n"
            if isinstance(json_data, dict):
                preview += f"对象，包含 {len(json_data)} 个键:\n"
                for key in list(json_data.keys())[:5]:
                    value_type = type(json_data[key]).__name__
                    preview += f"  - {key}: {value_type}\n"
                if len(json_data) > 5:
                    preview += "  ..."
            elif isinstance(json_data, list):
                preview += f"数组，包含 {len(json_data)} 个元素\n"
                if len(json_data) > 0:
                    preview += f"元素类型: {type(json_data[0]).__name__}"
            
            return {
                'content': content,
                'preview': preview
            }
        except json.JSONDecodeError as e:
            return {
                'content': f"JSON格式错误: {str(e)}",
                'preview': "JSON解析失败"
            }
    
    @staticmethod
    def process_xml_file(file) -> Dict[str, str]:
        """处理XML文件"""
        try:
            content = file.read().decode('utf-8')
            root = ET.fromstring(content)
            
            preview = f"XML文件结构:\n"
            preview += f"根元素: {root.tag}\n"
            
            # 显示子元素
            children = list(root)
            if children:
                preview += f"子元素 ({len(children)}):\n"
                for child in children[:5]:
                    preview += f"  - {child.tag}\n"
                if len(children) > 5:
                    preview += "  ..."
            
            return {
                'content': content,
                'preview': preview
            }
        except ET.ParseError as e:
            return {
                'content': f"XML解析错误: {str(e)}",
                'preview': "XML格式错误"
            }
    
    @staticmethod
    def process_excel_file(file) -> Dict[str, str]:
        """处理Excel文件"""
        if not OFFICE_SUPPORT:
            return {
                'content': "缺少Office文件支持库",
                'preview': "请安装 openpyxl 包"
            }
        
        try:
            # 保存临时文件
            import tempfile
            with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as tmp_file:
                for chunk in file.chunks():
                    tmp_file.write(chunk)
                tmp_file_path = tmp_file.name
            
            try:
                workbook = load_workbook(tmp_file_path)
                sheets = workbook.sheetnames
                
                preview = f"Excel文件 ({len(sheets)} 个工作表):\n"
                
                for sheet_name in sheets[:3]:  # 只处理前3个工作表
                    sheet = workbook[sheet_name]
                    max_row = min(sheet.max_row, 10)  # 最多读10行
                    max_col = min(sheet.max_column, 5)  # 最多读5列
                    
                    preview += f"\n工作表: {sheet_name}\n"
                    for row in range(1, max_row + 1):
                        row_data = []
                        for col in range(1, max_col + 1):
                            cell_value = sheet.cell(row=row, column=col).value
                            row_data.append(str(cell_value) if cell_value is not None else "")
                        preview += f"  {' | '.join(row_data)}\n"
                    
                    if sheet.max_row > 10:
                        preview += f"  ... (还有 {sheet.max_row - 10} 行)\n"
                
                # 转换为CSV格式作为内容
                first_sheet = workbook[sheets[0]]
                csv_content = []
                for row in first_sheet.iter_rows(values_only=True):
                    csv_content.append(','.join([str(cell) if cell is not None else "" for cell in row]))
                
                return {
                    'content': '\n'.join(csv_content),
                    'preview': preview,
                    'sheets': len(sheets)
                }
                
            finally:
                os.unlink(tmp_file_path)
                
        except Exception as e:
            return {
                'content': f"Excel处理错误: {str(e)}",
                'preview': "Excel文件处理失败"
            }
    
    @staticmethod
    def process_word_file(file) -> Dict[str, str]:
        """处理Word文档"""
        if not OFFICE_SUPPORT:
            return {
                'content': "缺少Office文件支持库",
                'preview': "请安装 python-docx 包"
            }
        
        try:
            # 保存临时文件
            import tempfile
            with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as tmp_file:
                for chunk in file.chunks():
                    tmp_file.write(chunk)
                tmp_file_path = tmp_file.name
            
            try:
                doc = Document(tmp_file_path)
                
                # 提取文本内容
                content = []
                for paragraph in doc.paragraphs:
                    if paragraph.text.strip():
                        content.append(paragraph.text)
                
                full_text = '\n'.join(content)
                preview = full_text[:1000] + "..." if len(full_text) > 1000 else full_text
                
                return {
                    'content': full_text,
                    'preview': f"Word文档 ({len(doc.paragraphs)} 段落)\n\n{preview}",
                    'paragraphs': len(doc.paragraphs)
                }
                
            finally:
                os.unlink(tmp_file_path)
                
        except Exception as e:
            return {
                'content': f"Word处理错误: {str(e)}",
                'preview': "Word文档处理失败"
            }
    
    @staticmethod
    def process_powerpoint_file(file) -> Dict[str, str]:
        """处理PowerPoint文件"""
        if not OFFICE_SUPPORT:
            return {
                'content': "缺少Office文件支持库",
                'preview': "请安装 python-pptx 包"
            }
        
        try:
            # 保存临时文件
            import tempfile
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
                for chunk in file.chunks():
                    tmp_file.write(chunk)
                tmp_file_path = tmp_file.name
            
            try:
                prs = Presentation(tmp_file_path)
                
                content = []
                preview = f"PowerPoint演示文稿 ({len(prs.slides)} 张幻灯片):\n\n"
                
                for i, slide in enumerate(prs.slides[:5]):  # 只处理前5张幻灯片
                    slide_text = []
                    preview += f"幻灯片 {i+1}:\n"
                    
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text)
                            preview += f"  {shape.text[:100]}...\n" if len(shape.text) > 100 else f"  {shape.text}\n"
                    
                    content.extend(slide_text)
                    preview += "\n"
                
                if len(prs.slides) > 5:
                    preview += f"... (还有 {len(prs.slides) - 5} 张幻灯片)"
                
                return {
                    'content': '\n'.join(content),
                    'preview': preview,
                    'slides': len(prs.slides)
                }
                
            finally:
                os.unlink(tmp_file_path)
                
        except Exception as e:
            return {
                'content': f"PowerPoint处理错误: {str(e)}",
                'preview': "PowerPoint文件处理失败"
            }
    
    @staticmethod
    def process_code_file(file) -> Dict[str, str]:
        """处理代码文件"""
        try:
            content = file.read().decode('utf-8')
            lines = content.split('\n')
            preview = f"代码文件 ({len(lines)} 行):\n\n"
            
            # 显示前20行作为预览
            for i, line in enumerate(lines[:20]):
                preview += f"{i+1:3d}: {line}\n"
            
            if len(lines) > 20:
                preview += f"\n... (还有 {len(lines) - 20} 行)"
            
            return {
                'content': content,
                'preview': preview,
                'lines': len(lines)
            }
        except UnicodeDecodeError:
            return {
                'content': "无法解码代码文件",
                'preview': "编码错误"
            }


def process_uploaded_files(files: List) -> List[Dict[str, Any]]:
    """
    批量处理上传的文件
    
    Args:
        files: 上传的文件列表
        
    Returns:
        处理后的文件信息列表
    """
    processed_files = []
    
    for file in files:
        try:
            file_info = FileProcessor.process_file(file)
            processed_files.append(file_info)
        except Exception as e:
            logger.error(f"处理文件失败: {str(e)}")
            processed_files.append({
                'filename': getattr(file, 'name', 'unknown'),
                'content': f"文件处理失败: {str(e)}",
                'preview': "处理错误",
                'file_type': 'error'
            })
    
    return processed_files